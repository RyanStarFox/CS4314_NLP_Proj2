import os
from typing import List, Dict, Optional
import docx2txt
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import base64
import io
from openai import OpenAI
from config import DATA_DIR, OPENAI_API_KEY, OPENAI_API_BASE, ENABLE_IMAGE_CAPTIONING, IMAGE_CAPTION_MODEL, VL_API_KEY, VL_API_BASE, get_openai_client

class DocumentLoader:
    def __init__(
        self,
        data_dir: str = DATA_DIR,
    ):
        self.data_dir = data_dir
        self.supported_formats = [".pdf", ".pptx", ".docx", ".txt", ".md"]
        
        # 初始化用于 Image Captioning 的客户端
        if ENABLE_IMAGE_CAPTIONING:
            self.client = get_openai_client(api_key=VL_API_KEY, base_url=VL_API_BASE)

    def _generate_image_caption(self, image_bytes: bytes, source_info: str = "") -> str:
        """调用视觉大模型生成图片描述"""
        if not ENABLE_IMAGE_CAPTIONING:
            return ""
            
        try:
            base64_image = base64.b64encode(image_bytes).decode('utf-8')
            
            response = self.client.chat.completions.create(
                model=IMAGE_CAPTION_MODEL,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "请详细描述这张图片的内容。如果是图表，请解释其数据趋势；如果是公式，请转写为LaTeX；如果是流程图，请说明步骤。请只输出描述内容，不要废话。"},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{base64_image}"
                                },
                            },
                        ],
                    }
                ],
                max_tokens=300,
            )
            caption = response.choices[0].message.content
            print(f" [图片理解] 已生成描述 ({source_info}): {caption[:30]}...")
            return f"\n\n[图片描述 ({source_info})]: {caption}\n\n"
        except Exception as e:
            print(f" [图片理解] 生成失败: {e}")
            return ""
    
    def load_pdf(self, file_path: str) -> List[Dict]:
        """加载PDF文件，按页返回内容 (包含图片理解)"""
        pages = []
        try:
            # 尝试使用 PyMuPDF (fitz) 因为它提取图片更方便
            import fitz 
            doc = fitz.open(file_path)
            
            for page_num, page in enumerate(doc, start=1):
                text = page.get_text()
                
                # 提取并处理图片
                image_descriptions = ""
                if ENABLE_IMAGE_CAPTIONING:
                    image_list = page.get_images(full=True)
                    for img_index, img in enumerate(image_list):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        
                        # 忽略过小的图标或装饰性图片 (例如小于 5KB)
                        if len(image_bytes) < 5 * 1024:
                            continue
                            
                        desc = self._generate_image_caption(image_bytes, source_info=f"第{page_num}页 图片{img_index+1}")
                        image_descriptions += desc

                formatted_text = f"--- 第 {page_num} 页 ---\n{text}\n{image_descriptions}\n"
                pages.append({"text": formatted_text})
                
        except ImportError:
            print("Warning: PyMuPDF (fitz) module not found. Falling back to PyPDF2. Image context will be unavailable.")
            # Fallback to PyPDF2
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                formatted_text = f"--- 第 {page_num} 页 ---\n{text}\n"
                pages.append({"text": formatted_text})
                
        except Exception as e:
            print(f"Error loading PDF: {e}")
            
        return pages

    def load_pptx(self, file_path: str) -> List[Dict]:
        """加载PPT文件，按幻灯片返回内容 (包含图片理解)"""
        slides = []
        presentation = Presentation(file_path)
        
        for slide_num, slide in enumerate(presentation.slides, start=1):
            text_parts = []
            image_descriptions = ""
            
            for shape in slide.shapes:
                # 提取文本
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)
                
                # 提取图片
                if ENABLE_IMAGE_CAPTIONING:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image_bytes = shape.image.blob
                             # 忽略过小的图片
                            if len(image_bytes) < 5 * 1024:
                                continue
                            desc = self._generate_image_caption(image_bytes, source_info=f"幻灯片{slide_num}")
                            image_descriptions += desc
                        except Exception as e:
                            print(f"PPT图片提取失败: {e}")
            
            slide_text = "\n".join(text_parts)
            formatted_text = f"--- 幻灯片 {slide_num} ---\n{slide_text}\n{image_descriptions}\n"
            slides.append({"text": formatted_text})
        
        return slides

    def load_docx(self, file_path: str) -> str:
        """加载DOCX文件
        TODO: 实现DOCX文件加载
        要求：
        1. 使用docx2txt读取DOCX文件
        2. 返回文本内容
        """
        text = docx2txt.process(file_path)
        return text

    def load_txt(self, file_path: str) -> str:
        """加载TXT文件
        TODO: 实现TXT文件加载
        要求：
        1. 使用open读取TXT文件（注意使用encoding="utf-8"）
        2. 返回文本内容
        """
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        return text

    def load_markdown(self, file_path: str) -> str:
        """加载Markdown文件
        TODO: 实现Markdown文件加载
        要求：
        1. 使用open读取Markdown文件（注意使用encoding="utf-8"）
        2. 返回文本内容
        """
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        return text

    def load_document(self, file_path: str) -> List[Dict[str, str]]:
        """加载单个文档，PDF和PPT按页/幻灯片分割，返回文档块列表"""
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        documents = []

        if ext == ".pdf":
            pages = self.load_pdf(file_path)
            for page_idx, page_data in enumerate(pages, 1):
                documents.append(
                    {
                        "content": page_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": page_idx,
                    }
                )
        elif ext == ".pptx":
            slides = self.load_pptx(file_path)
            for slide_idx, slide_data in enumerate(slides, 1):
                documents.append(
                    {
                        "content": slide_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": slide_idx,
                    }
                )
        elif ext == ".docx":
            content = self.load_docx(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        elif ext == ".txt":
            content = self.load_txt(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        elif ext == ".md":
            content = self.load_markdown(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        else:
            print(f"不支持的文件格式: {ext}")

        return documents

    def load_all_documents(self) -> List[Dict[str, str]]:
        """加载数据目录下的所有文档"""
        if not os.path.exists(self.data_dir):
            print(f"数据目录不存在: {self.data_dir}")
            return None

        documents = []

        for root, dirs, files in os.walk(self.data_dir):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in self.supported_formats:
                    file_path = os.path.join(root, file)
                    print(f"正在加载: {file_path}")
                    doc_chunks = self.load_document(file_path)
                    if doc_chunks:
                        documents.extend(doc_chunks)

        return documents
